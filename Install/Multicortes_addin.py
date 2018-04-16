#-*- coding: utf-8 -*-
import arcpy
import pythonaddins
import os,xlwt,xlrd
from xlrd import open_workbook
import arcgisscripting as script_tools
import subprocess
import exceptions
import inspect
import locale
import shutil
import xlsxwriter
import easygui
from tkFileDialog import askopenfilename
from tkFileDialog import asksaveasfilename
from tkFileDialog import askdirectory
import pptx
import pptx.util
import glob
import time
locale.setlocale(locale.LC_ALL,  'C')
import xlutils
from xlutils.copy import copy
from pptx.util import Inches,Pt,Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.lang import MSO_LANGUAGE_ID

import Tkinter, Tkconstants, tkFileDialog
from Tkinter import *

reload(sys) # cargue la codificación utf8
sys.setdefaultencoding('UTF8')

if not hasattr(sys, 'argv'):
    sys.argv = ['']

#=========Funciones Auxiliares=====================#
def getPythonPath():
    pydir = sys.exec_prefix
    pyexe = os.path.join(pydir, "python.exe")
    if os.path.exists(pyexe):
        return pyexe
    else:
        raise RuntimeError("python.exe no se encuentra instalado en {0}".format(pydir))

def directorioyArchivo ():
    archivo=inspect.getfile(inspect.currentframe()) # script filename
    directorio=os.path.dirname(os.path.abspath(inspect.getfile(inspect.currentframe()))) # script directory
    return archivo, directorio

# ------------------------------------------------------------

#=========Validación de requerimientos=====================#

pyexe = getPythonPath()
verPython64=pyexe
verPythonfinal=verPython64
# ------------------------------------------------------------


ruta=""
verPython32=verPython64
verPythonDir=verPython64.replace("\\python.exe","")
verPythonDir_32 =verPythonDir
ArcVersion=verPythonDir_32.replace("C:\Python27\ArcGIS","")
##archivo, directorio = directorioyArchivo()
home =os.path.expanduser("~")
directorio=home+"\\"+"Documents\ArcGIS\AddIns\Desktop"+str(arcpy.GetInstallInfo()['Version'])[:4]
directorio_raiz = directorio
script_clip=directorio+"\\clipx64_aux.py"
script_identity=directorio+"\\identityx64_aux.py"

def capstringTK(mensaje):
    try:
        sys.argv
    except:
        sys.argv=['']
    top = Tkinter.Tk()
    L1 =  Label(top, text=mensaje)
    L1.pack( side = LEFT)
    N1 = Entry(top, bd =5, width=50)
    N1.pack(side = LEFT)
    N1.focus_set()
    def callback():
        global valor_string
        valor_string = str(N1.get().encode('utf-8'))
        top.destroy()
    B1 = Button(top, text="Aceptar", command=callback)
    B1.pack(side= RIGHT)
    top.mainloop()
    return valor_string

def creadirs(ruta_raiz,nombre): # crea los dierctorios de salida del programa
    if not os.path.exists(ruta_raiz+"\\%s"%(nombre)):
        os.makedirs(ruta_raiz+"\\%s"%(nombre))
    return ruta_raiz+"\\%s"%nombre

def listar_xls(path,lst):
      for file in os.listdir(path):
        if file.endswith(".xls") or file.endswith(".xlsx"):
            lst.append(os.path.join(path, file))
      return lst

def fusion(path,nombre_capa):
    wkbk = xlwt.Workbook()
    lst=[]
    xlsfiles = listar_xls(path,lst)
    outrow_idx = 0
    for f in xlsfiles:
        insheet = xlrd.open_workbook(f).sheets()[0]
        outsheet = wkbk.add_sheet(xlrd.open_workbook(f).sheets()[0].name)
        for row_idx in xrange(insheet.nrows):
            for col_idx in xrange(insheet.ncols):
                outsheet.write(outrow_idx, col_idx,
                insheet.cell_value(row_idx, col_idx))
            outrow_idx += 1
        outrow_idx = 0
    wkbk.save(os.path.join(path, r'%s.xls'%(nombre_capa)))
    del wkbk,xlsfiles,outsheet

def grafica(f):
    outrow_idx = 0
    rb = open_workbook(f)
    rs = rb.sheet_by_index(0)
    wb = copy(rb)
    ws = wb.get_sheet(0)
    filename, file_extension = os.path.splitext(f)
    new = filename + r'_graf.xlsx'
    workbook  = xlsxwriter.Workbook(new)
    #worksheet = workbook.add_worksheet(ws.name)
    for s in rb.sheets():
        print 'Sheet:',s.name
        ws = wb.get_sheet(s.number)
        worksheet = workbook.add_worksheet(s.name)
        for row_idx in xrange(s.nrows):
            for col_idx in xrange(s.ncols):
                worksheet.write(outrow_idx, col_idx,s.cell_value(row_idx, col_idx))
            outrow_idx += 1
        outrow_idx = 0
        values = {0:"No Apto",1:"Baja",2:"Media",3:"Alta",4:"Na 4",5:"Na 5",6:"Na 6",7:"Na 7",8:"Exclusiones"}
        worksheet.write(0, 6, 'Aptitud')
        chart = workbook.add_chart({'type': 'pie'})
        chart.set_legend({'none': True})
        chart.set_size({'width': 335, 'height': 142})
        chart.set_style(18)
        chart.set_chartarea({
            'border': {'color': '#D9D9D9', 'width': 0.75}
        })
        cat = s.name+'!$G$2:$G$7'
        val = s.name+'!$F$2:$F$7'
        chart.add_series({
        #'name': s.name[0:-4],
        'categories': cat,
        'values':     val,
        'border': {'color': 'black'},
        'data_labels': {'value': True, 'leader_lines': True, 'font': {'name': 'Century Gothic', 'size': 11 ,'bold': True}, 'num_format': '##"%"' },
        'points': [
            {'fill': {'color': '#E1E1E1'}},
            {'fill': {'color': '#E9FFBE'}},
            {'fill': {'color': '#38D400'}},
            {'fill': {'color': '#266600'}},
            {'fill': {'color': '#828282'}},
            ],
        })
        #chart.set_title({'name': 'Aptitud %s'%(s.name[2:-4].replace("_", " "))})
        for i in range(1,8):
##                    print i
            try:
                a1 = int(s.cell_value(rowx=i, colx=1))
                worksheet.write(i, 6, values[a1])
                worksheet.insert_chart('H1', chart)
                del chart
            except:
                continue
    workbook.close()

def calc_area(capas):
    datos_area=[]
    datos_area_totales=[]
    datos_capa=[]
    for capa in capas:
        desc = arcpy.Describe(capa)
        with arcpy.da.SearchCursor(capa,"SHAPE@AREA") as cursor:
            for fila in cursor:
                datos_area.append(fila[0])
        datos_area_totales.append(sum(datos_area))
        datos_area = []
        datos_capa.append(desc.name)
    for x in xrange(len(datos_area_totales)):
        print "la capa {} tiene un area de {:.3f} en m2 y {} en ha".format(datos_capa[x],datos_area_totales[x],datos_area_totales[x]/10000)

def calc_ruta(capas):
    for capa in capas:
        desc = arcpy.Describe(capa)
        print "la capa {} tiene como ruta: {}".format(desc.name,desc.catalogpath)

def calc_coordenadas(capas):
    for capa in capas:
        desc = arcpy.Describe(capa)
        print "la capa {} tiene como sistema de coordenadas: {}".format(desc.name,desc.spatialreference.name)


def calc_porcentaje(tabla,campo_estadisticas):
    aRangosUnicos = []
    aRangos = []
    aValores = []
    aPorcentajes = []
    dicSumatorias = {}
    tabla_cursor=tabla
    desc=arcpy.Describe(tabla)
    with arcpy.da.SearchCursor(tabla_cursor,[campo_estadisticas,"SUM_Shape_Area"]) as cursor1:

    #Recorrido del Primer Cursor
        try:
            for filas in cursor1:
                try:
                    rango = filas[0]
                except:
                    rango = filas[0]
            	valor = filas[1]
            	aRangos.append(rango)
            	aValores.append(valor)
            	if rango not in aRangosUnicos:
            		aRangosUnicos.append(rango)
        except:
            nombre_tabla= desc.name
            for filas in cursor1:
                try:
                    rango = filas[0]
                except:
                    rango = filas[0]
            	valor = filas[1]
            	aRangos.append(rango)
            	aValores.append(valor)
            	if rango not in aRangosUnicos:
            		aRangosUnicos.append(rango)

    #Valores Sumatoria
    for k in aRangosUnicos:
        aSumatorias = []
    try:
            try:
                cursorAux = arcpy.SearchCursor("%s"%(tabla_cursor)," %s = '"%(campo_estadisticas)+k+"'",fields ="SUM_Shape_Area")
            except:
                cursorAux = arcpy.SearchCursor("%s"%(tabla_cursor)," %s = '"%(campo_estadisticas_tentativo)+k+"'",fields ="SUM_Shape_Area")

    except:
            try:
                cursorAux = arcpy.SearchCursor("%s"%(tabla_cursor)," %s = "%(campo_estadisticas)+str(k),fields ="SUM_Shape_Area")
            except:
                cursorAux = arcpy.SearchCursor("%s"%(tabla_cursor)," %s = "%(campo_estadisticas_tentativo)+str(k),fields ="SUM_Shape_Area")

    for f in cursorAux:
    		aSumatorias.append(f.getValue("SUM_Shape_Area"))
    suma = sum(aSumatorias)
    dicSumatorias[k]=suma

    #Calculo Porcentajes
    for i in range(len(aRangos)):
    	porcentaje = (float(aValores[i])/sum(aValores))*100
    	aPorcentajes.append(porcentaje)

    #Agregar columna a la tabla
    arcpy.AddField_management("%s"%(tabla_cursor),"SUM_Area_HA","DOUBLE")
    arcpy.AddField_management("%s"%(tabla_cursor),"Porcentaje","DOUBLE")
    p = 0
    tablaObjetivo = arcpy.UpdateCursor("%s"%(tabla_cursor))
    for row in tablaObjetivo:
        row.Porcentaje = aPorcentajes[p]
        row.SUM_Area_HA = row.SUM_Shape_Area/10000
        tablaObjetivo.updateRow(row)
        p = p+1


def ventana(arreglo):
    listOfOptions = arreglo
    GetUserInput(listOfOptions, True).getInput()



def toExcel():
    mxd= arcpy.mapping.MapDocument("CURRENT")
    tablas=[tabla.name for tabla in arcpy.mapping.ListTableViews(mxd)]
    global ruta
    ruta=pythonaddins.OpenDialog("Folder de almacenamiento",)
    ventana(tablas)




def calc_estadisticas(capas):
    for capa in capas:
        if "--" in capa.name:
            nombre_capa,campo=capa.name.split("--")
            tabla=arcpy.Statistics_analysis(in_table=capa, out_table="in_memory//%s_EST"%(nombre_capa), statistics_fields="Shape_Area SUM", case_field="%s"%(campo)).getOutput(0)
            calc_porcentaje(tabla,campo)
        else:
            print "la capa {} no tiene definido ningun campo de estadisticas".format(capa)



class GetUserInput(object):
    selection = None

    def __init__(self, options, multiple):
        self.master = Tk()

        self.master.title("Choose from list")
        self.scrollbar = Scrollbar(self.master, orient=VERTICAL)####

        self.listbox = Listbox(self.master, yscrollcommand=self.scrollbar.set,selectmode=EXTENDED if multiple else SINGLE, width=40, height=20)
        self.scrollbar.config(command=self.listbox.yview)##
        self.scrollbar.pack(side=RIGHT, fill=Y)##
        for option in options:
            self.listbox.insert(0, option)
##        self.listbox.pack()
        self.listbox.pack(side=LEFT, fill=BOTH, expand=1)##

        b = Button(self.master, command=self.callback, text="Convertir")
        b.pack()

        self.master.mainloop()

    def callback(self):
        self.selection = self.listbox.selection_get()
        capas= self.selection.split("\n")
        print capas
        for capa in capas:
            if "--" in capa:
                    nombre_capa,campo=capa.split("--")
                    arcpy.TableToExcel_conversion(capa, "%s//%s.xls"%(ruta,nombre_capa))
            else:
                arcpy.TableToExcel_conversion(capa, "%s//%s.xls"%(ruta,capa))
        self.master.destroy()

    def getInput(self):
        return self.selection


def mExport():
    import pythonaddins, time

    try:
        orgin= r'' + easygui.fileopenbox("Plantilla","Abrir archivo",".pptx",["*.pptx", ["*.ppt", "*.pptx", "Archivo de PowerPoint"]  ])
    except:
        orgin= r'' + askopenfilename(title="Seleccione un archivo",initialdir="C:",filetypes=[("Powerpoint",("*.ppt","*.pptx"))])
    try:
        out= r'' + easygui.filesavebox("Nombre el archivo","Guardar","out.pptx",["*.pptx", ["*.ppt", "*.pptx", "PowerPoint Files"]  ])
    except:
        out= r''+ asksaveasfilename(title="Nombre el archivo",initialdir="C:",initialfile="out.pptx",filetypes=[("Powerpoint",("*.ppt","*.pptx"))])
    try:
        PNGPath = r'' + easygui.diropenbox("Save Folder","Imagenes Generadas",".")
    except:
        PNGPath =r'' + askdirectory (title="Carpeta de almacenamiento",initialdir=".")
    formatoSalida =  easygui.choicebox("Seleccione formato de Salida","Formato",["PNG","JPG","PDF"])
    resolucionSalida = easygui.integerbox('Resolución 150 - 600 px',' Resolución Salida',300,150,601)
    t_inicio=time.clock()
    mxd = arcpy.mapping.MapDocument("CURRENT")
    df = arcpy.mapping.ListDataFrames(mxd, '')[0]
    lyrList = pythonaddins.GetSelectedTOCLayerOrDataFrame()
    prs = pptx.Presentation(orgin)
    pic_left  = int(prs.slide_width - prs.slide_height)
    pic_top   = int(0)
    pic_width = int(prs.slide_height)
    pic_height = int(prs.slide_height)
    altoHoja = mxd.pageSize.height
    anchoHoja = mxd.pageSize.width

    alHojaPulga = altoHoja*0.39370
    anHojaPulga = anchoHoja*0.39370

    for lyr in arcpy.mapping.ListLayers(mxd, '', df):
        for layer in lyrList:
            if lyr.name == layer.name:
                lyr.visible = False
    arcpy.RefreshActiveView()

    if formatoSalida == "PNG":
        for lyr in arcpy.mapping.ListLayers(mxd, '', df):
            for layer in lyrList:
                if lyr.name == layer.name:
                    #pythonaddins.MessageBox("Exportando {} to PNG ".format(lyr.name),"Wait...",)
                    print "Exportando {} to PNG ".format(lyr.name)
                    lyr.visible = True
                    arcpy.RefreshActiveView()
                    # print mxd, PNGPath+"\\" + lyr.name + ".png","",anHojaPulga,alHojaPulga,int(resolucionSalida)
                    try:
                        arcpy.mapping.ExportToPNG(mxd, PNGPath+"\\" + lyr.name + ".png","",anHojaPulga,alHojaPulga,int(resolucionSalida))
                    except:
                        try:
                            arcpy.RefreshActiveView()
                            arcpy.mapping.ExportToPNG(mxd, PNGPath+"\\" + lyr.name + ".png","",anHojaPulga,alHojaPulga,int(resolucionSalida))
                        except:
                            pythonaddins.MessageBox("AttributeError: PageLayoutObject \n\nPlease restart ArcMap and try again , If error persist call to ESRI support","ArcGIS Fatal Error",7)
                            raise SystemExit(0)
                    lyr.visible = False
                    #Agrega imagen a pptx
                    g = glob.glob(PNGPath+"\\" + lyr.name + ".png")[0]
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    tb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1.51), Cm(4.7), Cm(13) , Cm(2.31))
                    tb.fill.background()
                    tb.line.fill.background()
                    p = tb.text_frame.add_paragraph()
                    p.alignment = PP_ALIGN.CENTER
                    p.text = lyr.name
                    p.font.language_id=MSO_LANGUAGE_ID.SPANISH_COLOMBIA
                    p.font.size = Pt(24)
                    p.font.color.rgb = RGBColor(21, 151, 65)
                    p.font.bold = True
                    p.font.name = 'Century Gothic'
                    pic = slide.shapes.add_picture(g, pic_left, pic_top, pic_width, pic_height)
                    try:
                        top = 6.8
                        for i in lyr.symbology.classBreakLabels:
                            tc = slide.shapes.add_shape(MSO_SHAPE.ROUND_2_SAME_RECTANGLE, Cm(5.45), Cm(top) , Cm(1.16) , Cm(4.16))
                            tc.rotation = -90.0
                            tc.line.fill.background()
                            tc.text_frame.margin_bottom = Cm(0.13)
                            tc.text_frame.margin_left = Cm(0.25)
                            tc.text_frame.margin_right = Cm(0.25)
                            tc.text_frame.margin_top = Cm(0.13)
                            q = tc.text_frame.add_paragraph()
                            q.alignment = PP_ALIGN.CENTER
                            q.text = str(i)
                            q.font.language_id=MSO_LANGUAGE_ID.SPANISH_COLOMBIA
                            q.font.size = Pt(14)
                            q.font.color.rgb = RGBColor(0, 0, 0)
                            q.font.bold = True
                            q.font.name = 'Century Gothic'
                            top += 1.5
                    except:
                        print "Group Layer no has Symbology"
                    del g
        prs.save(out)

    elif formatoSalida == "PDF":
        easygui.msgbox("PDF format not supported by PowerPoint integration!!! \n Only Export images")
        for lyr in arcpy.mapping.ListLayers(mxd, '', df):
            for layer in lyrList:
                if lyr.name == layer.name:
                    #pythonaddins.MessageBox("Exportando {} to PDF ".format(lyr.name),"Wait...")
                    print "Exportando {} to PDF ".format(lyr.name)
                    lyr.visible = True
                    arcpy.RefreshActiveView()
                    try:
                        arcpy.mapping.ExportToPDF(mxd, PNGPath+"\\" + lyr.name + ".pdf","",anHojaPulga,alHojaPulga,int(resolucionSalida))
                    except:
                        try:
                            arcpy.RefreshActiveView()
                            arcpy.mapping.ExportToPDF(mxd, PNGPath+"\\" + lyr.name + ".pdf","PAGE_LAYOUT",anHojaPulga,alHojaPulga,int(resolucionSalida))
                        except:
                            pythonaddins.MessageBox("AttributeError: PageLayoutObject \n\nPlease restart ArcMap and try again , If error persist call to ESRI support","ArcGIS Fatal Error",7)
                            raise SystemExit(0)
                    lyr.visible = False
                    try:
                        #Agrega imagen a pptx
                        g = glob.glob(PNGPath+"\\" + lyr.name + ".pdf")[0]
                        print g
                        slide = prs.slides.add_slide(prs.slide_layouts[6])
                        tb = slide.shapes.add_textbox(Cm(1.51), Cm(4.7), Cm(13) , Cm(2.31))
                        p = tb.text_frame.add_paragraph()
                        p.alignment = PP_ALIGN.CENTER
                        p.text = lyr.name
                        p.font.language_id=MSO_LANGUAGE_ID.SPANISH_COLOMBIA
                        p.font.size = Pt(24)
                        p.font.color.rgb = RGBColor(21, 151, 65)
                        p.font.bold = True
                        p.font.name = 'Century Gothic'
                        pic = slide.shapes.add_picture(g, pic_left, pic_top, pic_width, pic_height)
                        try:
                            tc = slide.shapes.add_textbox(Cm(1.51), Cm(6.7), Cm(13) , Cm(2.31))
                            q = tc.text_frame.add_paragraph()
                            q.text = str(lyr.symbology.classBreakLabels)
                            q.font.language_id=MSO_LANGUAGE_ID.SPANISH_COLOMBIA
                            q.font.size = Pt(14)
                            q.font.color.rgb = RGBColor(21, 151, 65)
                            q.font.bold = True
                            q.font.name = 'Century Gothic'
                        except:
                            print "Group Layer no has Symbology"
                    except:
                        print "PDF is not supported by PowerPoint insert"
                    del g
        prs.save(out)


    else:
        for lyr in arcpy.mapping.ListLayers(mxd, '', df):
            for layer in lyrList:
                if lyr.name == layer.name:
                    #pythonaddins.MessageBox("Exportando {} to JPG ".format(lyr.name),"Wait...")
                    print "Exportando {} to JPG ".format(lyr.name)
                    lyr.visible = True
                    arcpy.RefreshActiveView()
                    try:
                        arcpy.mapping.ExportToJPEG(mxd, PNGPath+"\\" + lyr.name + ".jpg","",anHojaPulga,alHojaPulga,int(resolucionSalida))
                    except:
                        try:
                            arcpy.RefreshActiveView()
                            arcpy.mapping.ExportToJPEG(mxd, PNGPath+"\\" + lyr.name + ".jpg","PAGE_LAYOUT",anHojaPulga,alHojaPulga,int(resolucionSalida))
                        except:
                            pythonaddins.MessageBox("AttributeError: PageLayoutObject \n\nPlease restart ArcMap and try again , If error persist call to ESRI support","ArcGIS Fatal Error",7)
                            raise SystemExit(0)
                    lyr.visible = False
                    #Agrega imagen a pptx
                    #Agrega imagen a pptx
                    g = glob.glob(PNGPath+"\\" + lyr.name + ".jpg")[0]
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    tb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1.51), Cm(4.7), Cm(13) , Cm(2.31))
                    tb.fill.background()
                    tb.line.fill.background()
                    p = tb.text_frame.add_paragraph()
                    p.alignment = PP_ALIGN.CENTER
                    p.text = lyr.name
                    p.font.language_id=MSO_LANGUAGE_ID.SPANISH_COLOMBIA
                    p.font.size = Pt(24)
                    p.font.color.rgb = RGBColor(21, 151, 65)
                    p.font.bold = True
                    p.font.name = 'Century Gothic'
                    pic = slide.shapes.add_picture(g, pic_left, pic_top, pic_width, pic_height)
                    try:
                        top = 6.8
                        for i in lyr.symbology.classBreakLabels:
                            tc = slide.shapes.add_shape(MSO_SHAPE.ROUND_2_SAME_RECTANGLE, Cm(5.45), Cm(top) , Cm(1.16) , Cm(4.16))
                            tc.rotation = -90.0
                            tc.line.fill.background()
                            tc.text_frame.margin_bottom = Cm(0.13)
                            tc.text_frame.margin_left = Cm(0.25)
                            tc.text_frame.margin_right = Cm(0.25)
                            tc.text_frame.margin_top = Cm(0.13)
                            q = tc.text_frame.add_paragraph()
                            q.alignment = PP_ALIGN.CENTER
                            q.text = str(i)
                            q.font.language_id=MSO_LANGUAGE_ID.SPANISH_COLOMBIA
                            q.font.size = Pt(14)
                            q.font.color.rgb = RGBColor(0, 0, 0)
                            q.font.bold = True
                            q.font.name = 'Century Gothic'
                            top += 1.5
                    except:
                        print "Group Layer no has Symbology"
                    del g
        prs.save(out)

    start = t_inicio
    end = time.clock()
    hours, rem = divmod(end-start, 3600)
    minutes, seconds = divmod(rem, 60)
    print("proceso Completado en {:0>2} H {:0>2} M {:05.2f} S.".format(int(hours),int(minutes),seconds))




class ButtonArea(object):
    """Implementation for UpraToolBar_InfoCapas.button (Button)"""
    def __init__(self):
        self.enabled = True
        self.checked = False
    def onClick(self):
        capas=pythonaddins.GetSelectedTOCLayerOrDataFrame()
        if capas is not None:
            if type(capas)!= list:
                arr =[]
                arr.append(capas)
                capas=arr
                calc_area(arr)
            else:
                if len(capas)>=1:
                    calc_area(capas)
                else:
                    print"###### Seleccione por lo menos una capa en el dataframe activo ######"
        else:
            print" ###### Seleccione por lo menos una capa en el dataframe activo ######"




class ButtonCoordenadas(object):
    """Implementation for UpraToolBar_Coordenadas.button (Button)"""
    def __init__(self):
        self.enabled = True
        self.checked = False
    def onClick(self):
        capas=pythonaddins.GetSelectedTOCLayerOrDataFrame()
        if capas is not None:
            if type(capas)!= list:
                arr =[]
                arr.append(capas)
                capas=arr
                calc_coordenadas(arr)
            else:
                if len(capas)>=1:
                    calc_coordenadas(capas)
                else:
                    print "###### Seleccione por lo menos una capa en el dataframe activo ######"
        else:
            print "###### Seleccione por lo menos una capa en el dataframe activo ######"



class ButtonEstadisticas(object):
    """Implementation for UpraToolBar_Estadisticas.button (Button)"""
    def __init__(self):
        self.enabled = True
        self.checked = False
    def onClick(self):
        capas=pythonaddins.GetSelectedTOCLayerOrDataFrame()
        if not hasattr(sys, 'argv'):
            sys.argv  = ['']
        if capas is not None:
            if type(capas)!= list:
                arr =[]
                arr.append(capas)
                capas=arr
                calc_estadisticas(arr)
            else:
                if len(capas)>=1:
                    calc_estadisticas(capas)
                else:
                    print "###### Seleccione por lo menos una capa en el dataframe activo ######"
        else:
            print "###### Seleccione por lo menos una capa en el dataframe activo ######"

class ButtonMulticortes(object):
    """Implementation for UpraToolBar_Presentaciones.ButtonMulticortes (Button)"""
    def __init__(self):
        self.enabled = True
        self.checked = False
    def onClick(self):
################################## MENU DE EJECUCION ##########################################################################
        #ruta=capstringTK("Ruta de Salida: ")
        if not hasattr(sys, 'argv'):
            sys.argv  = ['']
        try:
            ruta=easygui.diropenbox("Ruta de Salida: ","Select File",'C:')
        except:
            ruta=pythonaddins.OpenDialog("Ruta de Salida: ")
        #tipo_analisis=capstringTK("Seleccione el tipo de análisis 1=Corte y Estadísticas , 2=Solo corte , 3=Extracción por Campo con Estadísticas  : ")
        msg = "Seleccione el tipo de análisis"
        choices = ["Nada","Corte y Estadísticas","Solo corte","Extracción por Campo con Estadísticas"]
        tipo_analisis= easygui.indexbox(msg, choices=choices)
        if tipo_analisis==0:
            easygui.msgbox('GRACIAS por usar Multicortes Addin!!!  :p')
            return
        print "Tipo análisis" + str(tipo_analisis)
        if tipo_analisis!=3:
            tipo_corte=easygui.indexbox("Seleccione el tipo de Operación",  choices=["Clip" , "Identity"])
            if tipo_analisis!=2:
                try:
                    campo_estadisticas_tentativo=easygui.enterbox('Campo de Estadísticas','Write Field','gridcode')
                except:
                    campo_estadisticas_tentativo=easygui.textbox('Campo de Estadísticas','Write Field',['gridcode'],['gridcode'])[0]
            #ejecucion=capstringTK(": 1=Si , 2=No ")
            ejecucion=easygui.ynbox('¿Desea Continuar con la Ejecución?', 'Seguir?', ('Sí', 'No'))
        else:
            try:
                campo_extraccion=easygui.enterbox('Campo de Extracción','Write Field','gridcode')
            except:
                easygui.textbox('Campo de Extracción','Write Field',['gridcode'],['gridcode'])[0]
            if tipo_analisis!=2:
                try:
                    campo_estadisticas_tentativo=easygui.enterbox('Campo de Estadísticas','Write Field','gridcode')
                except:
                    easygui.textbox('Campo de Estadísticas','Write Field',['Campo Estadísticas:'],['gridcode'])[0]
            ejecucion=easygui.ynbox('¿Desea Continuar con la Ejecución?', 'Seguir?', ('Sí', 'No'))
################################## MENU DE EJECUCION ##########################################################################
        if ejecucion:
            doc=arcpy.mapping.MapDocument("CURRENT")
            df_molde=arcpy.mapping.ListDataFrames(doc,"*molde")
            df_capas=arcpy.mapping.ListDataFrames(doc,"*capas")
            lista_molde=arcpy.mapping.ListLayers(doc,"*",df_molde[0])
            lista_capas=arcpy.mapping.ListLayers(doc,"*",df_capas[0])
            if tipo_analisis==3:
                tipo_campo_extract=[campo.type for campo in arcpy.Describe(lista_molde[0]).fields if campo.name in ('%s'%(campo_extraccion))][0]
            ruta_capas=creadirs(ruta,"Capas")
            if tipo_analisis!="2":
                ruta_excel=creadirs(ruta,"Tablas_Excel")
            arreglo_rutas_excel=[]

            incremento=(100/(len(lista_molde)))
            print incremento
            with pythonaddins.ProgressDialog as dialog:
             dialog.title = "Ejecutando Proceso"
             dialog.description = "Procesando."
             dialog.animation = "Spiral"
##             dialog.progress=0
             for molde in lista_molde: ###############################################################################################
                dialog.description = "Procesando." + str(molde.name.encode("utf8"))
                if tipo_analisis!=2:
                    ruta_excel_temp=creadirs(ruta_excel,molde.name)
                    arreglo_rutas_excel.append(ruta_excel_temp)
                    campo_estadisticas=campo_estadisticas_tentativo
                rfinal=arcpy.CreateFileGDB_management(ruta_capas, molde.name+".gdb").getOutput(0)
                for capa in lista_capas:
                    dialog.description = "Procesando." + str(molde.name.encode("utf8")) +'\n' +str(capa.name.encode("utf8"))
                    temparr=capa.name.split("--")

                    if len(temparr)!=1:
                        arreglo_nombre=capa.name.split("--")
                        nombre_tabla=arreglo_nombre[0]
                        campo_estadisticas=arreglo_nombre[1]
                    else:
                        nombre_tabla=capa.name



                    if tipo_analisis!=3:
                        if tipo_corte==0:
                            moldex=arcpy.Describe(molde).catalogpath.encode("utf8")
                            capax=arcpy.Describe(capa).catalogpath.encode("utf8")
                            print moldex,capax,rfinal.encode("utf8"),nombre_tabla.encode("utf8")
                            arcpy.arcpy.Clip_analysis(in_features=capax, clip_features=moldex, out_feature_class=rfinal.encode("utf8")+"\\"+nombre_tabla.encode("utf8"))
                            arcpy.MakeFeatureLayer_management(rfinal+"\\"+nombre_tabla,nombre_tabla+"_"+molde.name)



                        elif tipo_corte==1:
##                                arcpy.Identity_analysis(molde,capa,rfinal+"\\"+nombre_tabla)
                            moldex=arcpy.Describe(molde).catalogpath.encode("utf8")
                            capax=arcpy.Describe(capa).catalogpath.encode("utf8")
                            print moldex,capax,rfinal.encode("utf8"),nombre_tabla.encode("utf8")
                            arcpy.Identity_analysis(in_features=moldex,identity_features=capax, out_feature_class=rfinal.encode("utf8")+"\\"+nombre_tabla.encode("utf8"),join_attributes="ALL")
                            arcpy.MakeFeatureLayer_management(rfinal+"\\"+nombre_tabla,nombre_tabla+"_"+molde.name)
                    else:
                        valor_extraccion=[row[0] for row in arcpy.da.SearchCursor(molde, campo_extraccion)][0]
                        if tipo_campo_extract in ('String'):
                            query_extract = "%s = '%s'"%(campo_extraccion,valor_extraccion)
                            arcpy.Select_analysis(capa,rfinal+"\\"+nombre_tabla,query_extract)
                        else:
                            query_extract = "%s = %s"%(campo_extraccion,valor_extraccion)
                            arcpy.Select_analysis(capa,rfinal+"\\"+nombre_tabla,query_extract)

                    aRangosUnicos = []
                    aRangos = []
                    aValores = []
                    aPorcentajes = []
                    dicSumatorias = {}

                    if tipo_analisis!=2:
                        if tipo_analisis!=3:
                            try:
                                arcpy.Statistics_analysis(in_table="%s"%(nombre_tabla+"_"+molde.name), out_table="%s//%s_EST"%(rfinal,nombre_tabla), statistics_fields="Shape_Area SUM", case_field="%s"%(campo_estadisticas))
                            except:
                                arcpy.Statistics_analysis(in_table="%s"%(nombre_tabla+"_"+molde.name), out_table="%s//%s_EST"%(rfinal,nombre_tabla), statistics_fields="Shape_Area SUM", case_field="%s"%(campo_estadisticas_tentativo))

                        else:
                            try:
                                arcpy.Statistics_analysis(in_table="%s"%(rfinal+"\\"+nombre_tabla), out_table="%s//%s_EST"%(rfinal,nombre_tabla), statistics_fields="Shape_Area SUM", case_field="%s"%(campo_estadisticas))
                            except:
                                arcpy.Statistics_analysis(in_table="%s"%(rfinal+"\\"+nombre_tabla), out_table="%s//%s_EST"%(rfinal,nombre_tabla), statistics_fields="Shape_Area SUM", case_field="%s"%(campo_estadisticas_tentativo))

                        tabla_cursor=rfinal+"\\"+nombre_tabla+"_EST"
                        cursor1 = arcpy.SearchCursor("%s"%(tabla_cursor))
                        nombre_tabla=nombre_tabla+"_EST"
                        #Recorrido del Primer Cursor
                        try:
                            for filas in cursor1:
                                try:
                                    rango = filas.getValue(campo_estadisticas)
                                except:
                                    rango = filas.getValue(campo_estadisticas_tentativo)
                            	valor = filas.getValue("SUM_Shape_Area")
                            	aRangos.append(rango)
                            	aValores.append(valor)
                            	if rango not in aRangosUnicos:
                            		aRangosUnicos.append(rango)
                        except:
                            nombre_tabla=arreglo_nombre[0]
                            campo_estadisticas=arreglo_nombre[1]
                            for filas in cursor1:
                                try:
                                    rango = filas.getValue(campo_estadisticas)
                                except:
                                    rango = filas.getValue(campo_estadisticas_tentativo)
                            	valor = filas.getValue("SUM_Shape_Area")
                            	aRangos.append(rango)
                            	aValores.append(valor)
                            	if rango not in aRangosUnicos:
                            		aRangosUnicos.append(rango)

                        #Valores Sumatoria
                        for k in aRangosUnicos:
            	            aSumatorias = []
    ##                        arcpy.AddMessage(k)
                        try:
                                tabla_cursor=rfinal+"\\"+nombre_tabla
                                try:
                                    cursorAux = arcpy.SearchCursor("%s"%(tabla_cursor)," %s = '"%(campo_estadisticas)+k+"'",fields ="SUM_Shape_Area")
                                except:
                                    cursorAux = arcpy.SearchCursor("%s"%(tabla_cursor)," %s = '"%(campo_estadisticas_tentativo)+k+"'",fields ="SUM_Shape_Area")

                        except:
                                tabla_cursor=rfinal+"\\"+nombre_tabla
                                try:
                                    cursorAux = arcpy.SearchCursor("%s"%(tabla_cursor)," %s = "%(campo_estadisticas)+str(k),fields ="SUM_Shape_Area")
                                except:
                                    cursorAux = arcpy.SearchCursor("%s"%(tabla_cursor)," %s = "%(campo_estadisticas_tentativo)+str(k),fields ="SUM_Shape_Area")



                        for f in cursorAux:
                        		aSumatorias.append(f.getValue("SUM_Shape_Area"))
                        suma = sum(aSumatorias)
                        dicSumatorias[k]=suma

                        #Calculo Porcentajes
                        for i in range(len(aRangos)):
                        	porcentaje = (float(aValores[i])/sum(aValores))*100
                        	aPorcentajes.append(porcentaje)

                        #Agregar columna a la tabla
                        arcpy.AddField_management("%s"%(tabla_cursor),"SUM_Area_HA","DOUBLE")
                        arcpy.AddField_management("%s"%(tabla_cursor),"Porcentaje","DOUBLE")

                        p = 0
                        tablaObjetivo = arcpy.UpdateCursor("%s"%(tabla_cursor))
                        for row in tablaObjetivo:
                            row.Porcentaje = aPorcentajes[p]
                            row.SUM_Area_HA = row.SUM_Shape_Area/10000
                            tablaObjetivo.updateRow(row)
                            p = p+1

                        arcpy.TableToExcel_conversion("%s//%s"%(rfinal,nombre_tabla), "%s//%s.xls"%(ruta_excel_temp,nombre_tabla))
                        if dialog.cancelled:
                            raise Exception("Error, Verifique sus datos")

                incremento+=incremento
                dialog.progress += incremento
                if tipo_analisis!=2:
                    fusion(ruta_excel_temp,ruta_excel_temp.split("\\")[-1])
                    grafica(ruta_excel_temp+"\\"+ruta_excel_temp.split("\\")[-1]+".xls")

        else:
            pass # Fin de la ejecución el usuario finalizó la ejecución.

def multiQuery():
    expresion = easygui.textbox(msg="Escriba el query",title="Definición de Query")
    mxd = arcpy.mapping.MapDocument("CURRENT")
    layers = pythonaddins.GetSelectedTOCLayerOrDataFrame()
    if len(layers) >1:
        for lyr in layers:
            lyr.definitionQuery = expresion
    elif len(layers) ==1:
        layers.definitionQuery = expresion
    else:
        print "######### SELECCIONE UN LAYER EN EL DATAFRAME ACTIVO"########
    arcpy.RefreshActiveView
    arcpy.RefreshTOC()

def multiExport():
    formatoSalida =  easygui.choicebox("Seleccione formato de Salida","Formato",["FeatureClass","Tabla"])
    try:
        ruta_salida = r'' + easygui.diropenbox("Caperta/GDB de Salida","Archivos de Salida",".")
    except:
        ruta_salida =r'' + askdirectory (title="Caperta/GDB de Salida",initialdir=".")
    lista_capas = pythonaddins.GetSelectedTOCLayerOrDataFrame()
    if lista_capas is not None and lista_capas != []:
        if type(lista_capas)!= list:
            arr =[]
            arr.append(lista_capas)
            lista_capas=arr
            for capa in lista_capas:
                if formatoSalida =="FeatureClass":
                    arcpy.FeatureClassToFeatureClass_conversion(capa,ruta_salida,capa.name)
                else:
                    arcpy.CopyRows_management(capa,os.path.join(ruta_salida,capa.name))
        else:
            for capa in lista_capas:
                if formatoSalida =="FeatureClass":
                    arcpy.FeatureClassToFeatureClass_conversion(capa,ruta_salida,capa.name)
                else:
                    arcpy.CopyRows_management(capa,os.path.join(ruta_salida,capa.name))

class ButtonRuta(object):
    """Implementation for UpraToolBar_Ruta.button (Button)"""
    def __init__(self):
        self.enabled = True
        self.checked = False
    def onClick(self):
        capas=pythonaddins.GetSelectedTOCLayerOrDataFrame()
        if capas is not None:
            if type(capas)!= list:
                arr =[]
                arr.append(capas)
                capas=arr
                calc_ruta(arr)
            else:
                if len(capas)>=1:
                    calc_ruta(capas)
                else:
                    print "###### Seleccione por lo menos una capa en el dataframe activo ######"
        else:
            print "###### Seleccione por lo menos una capa en el dataframe activo ######"

class ButtonToExcel(object):
    """Implementation for UpraToolBar_ToExcel.button (Button)"""
    def __init__(self):
        self.enabled = True
        self.checked = False
    def onClick(self):
        toExcel()

class ButtonToPowerpoint(object):
    """Implementation for UpraToolBar_Presentaciones.ButtonToPowerpoint (Button)"""
    def __init__(self):
        self.enabled = True
        self.checked = False
    def onClick(self):
        mExport()

class ButtonHelp(object):
    """Implementation for Multicortes_Help.button (Button)"""
    def __init__(self):
        self.enabled = True
        self.checked = False
    def onClick(self):
        global subprocess
        subprocess.Popen("start chrome /new-tab https://github.com/UpraAnalisis/UPRA-Analisis-Tools#herramientas-para-presentaciones",shell = True)

class ButtonMultiExport(object):
    """Implementation for UpraToolBar_MultiExport.button (Button)"""
    def __init__(self):
        self.enabled = True
        self.checked = False
    def onClick(self):
        capas=pythonaddins.GetSelectedTOCLayerOrDataFrame()
        if capas is not None:
            if type(capas)!= list:
                arr =[]
                arr.append(capas)
                capas=arr
                multiExport()
            else:
                if len(capas)>=1:
                     multiExport()
                else:
                    print "###### Seleccione por lo menos una capa en el dataframe activo ######"
        else:
            print "###### Seleccione por lo menos una capa en el dataframe activo ######"
        del capas

class ButtonMultiQuery(object):
    """Implementation for Multicortes_MultiQuery.button (Button)"""
    def __init__(self):
        self.enabled = True
        self.checked = False
    def onClick(self):
        capas=pythonaddins.GetSelectedTOCLayerOrDataFrame()
        if capas is not None:
            if type(capas)!= list:
                arr =[]
                arr.append(capas)
                capas=arr
                multiQuery()
            else:
                if len(capas)>=1:
                     multiQuery()
                else:
                    print "###### Seleccione por lo menos una capa en el dataframe activo ######"
        else:
            print "###### Seleccione por lo menos una capa en el dataframe activo ######"